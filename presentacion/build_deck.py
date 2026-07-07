# -*- coding: utf-8 -*-
"""Genera una presentacion bonita para el Taller de Agentes de IA (2h)."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.oxml.ns import qn
import copy

# ---------------- Paleta ----------------
INDIGO   = RGBColor(0x1A, 0x16, 0x4B)   # fondo oscuro
INDIGO2  = RGBColor(0x2A, 0x22, 0x6E)   # panel oscuro
VIOLET   = RGBColor(0x6D, 0x28, 0xD9)   # morado principal
VIOLET_L = RGBColor(0x8B, 0x5C, 0xF6)
CYAN     = RGBColor(0x22, 0xD3, 0xEE)   # acento cian
CORAL    = RGBColor(0xFB, 0x71, 0x85)   # acento coral
GOLD     = RGBColor(0xFB, 0xBF, 0x24)
LIGHT    = RGBColor(0xF6, 0xF5, 0xFB)   # fondo claro
CARD     = RGBColor(0xFF, 0xFF, 0xFF)
INK      = RGBColor(0x1E, 0x1B, 0x2E)   # texto oscuro
MUTED    = RGBColor(0x6B, 0x66, 0x82)   # texto gris
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)

HFONT = "Segoe UI Semibold"
BFONT = "Segoe UI"

EMU_IN = 914400
SW, SH = 13.333, 7.5

prs = Presentation()
prs.slide_width  = Inches(SW)
prs.slide_height = Inches(SH)
BLANK = prs.slide_layouts[6]

# ---------------- helpers ----------------
def slide():
    return prs.slides.add_slide(BLANK)

def bg(s, color):
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = color

def rect(s, x, y, w, h, color, line=None, line_w=None, shape=MSO_SHAPE.RECTANGLE):
    sp = s.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    if color is None:
        sp.fill.background()
    else:
        sp.fill.solid(); sp.fill.fore_color.rgb = color
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line
        sp.line.width = Pt(line_w or 1)
    sp.shadow.inherit = False
    return sp

def grad(sp, c1, c2, angle=45):
    """Aplica gradiente lineal a la forma."""
    spPr = sp.fill._xPr
    # quitar fill existente
    for tag in ('a:noFill','a:solidFill','a:gradFill','a:blipFill','a:pattFill','a:grpFill'):
        e = spPr.find(qn(tag))
        if e is not None: spPr.remove(e)
    ln = spPr.find(qn('a:ln'))
    grad = spPr.makeelement(qn('a:gradFill'), {})
    lst = grad.makeelement(qn('a:gsLst'), {})
    for pos, col in ((0, c1), (100000, c2)):
        gs = grad.makeelement(qn('a:gs'), {'pos': str(pos)})
        clr = grad.makeelement(qn('a:srgbClr'), {'val': '%02X%02X%02X' % (col[0], col[1], col[2])})
        gs.append(clr); lst.append(gs)
    grad.append(lst)
    lin = grad.makeelement(qn('a:lin'), {'ang': str(int(angle*60000)), 'scaled': '1'})
    grad.append(lin)
    if ln is not None:
        ln.addprevious(grad)
    else:
        spPr.append(grad)
    return sp

def text(s, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
         space_after=6, line_spacing=1.0, wrap=True):
    """runs: list of paragraphs; each paragraph is list of (txt, size, color, bold, font, italic)."""
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = 0; tf.margin_right = 0
    tf.margin_top = 0; tf.margin_bottom = 0
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space_after)
        p.space_before = Pt(0)
        p.line_spacing = line_spacing
        for (txt, size, color, bold, font, italic) in para:
            r = p.add_run(); r.text = txt
            r.font.size = Pt(size); r.font.bold = bold; r.font.italic = italic
            r.font.name = font; r.font.color.rgb = color
    return tb

def R(txt, size, color, bold=False, font=BFONT, italic=False):
    return (txt, size, color, bold, font, italic)

def circle_icon(s, x, y, d, fill, emoji, esize=20, ecolor=WHITE):
    c = rect(s, x, y, d, d, fill, shape=MSO_SHAPE.OVAL)
    text(s, x, y-0.02, d, d, [[R(emoji, esize, ecolor, True, "Segoe UI Emoji")]],
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, space_after=0)
    return c

def page_num(s, n=None, dark=False):
    n = len(prs.slides._sldIdLst)
    text(s, SW-1.0, SH-0.5, 0.6, 0.3, [[R(str(n), 11, (CYAN if dark else MUTED))]],
         align=PP_ALIGN.RIGHT, space_after=0)

def footer(s, dark=False):
    col = RGBColor(0x9A,0x94,0xC4) if dark else MUTED
    text(s, 0.6, SH-0.5, 6, 0.3, [[R("Taller de Agentes de IA  ·  Microsoft", 10, col)]],
         align=PP_ALIGN.LEFT, space_after=0)

# ================= SLIDE 1: TITULO =================
s = slide(); bg(s, INDIGO)
# formas decorativas
c = rect(s, 9.4, -1.8, 5.5, 5.5, VIOLET, shape=MSO_SHAPE.OVAL); grad(c, (0x6D,0x28,0xD9),(0x22,0xD3,0xEE),35)
c2 = rect(s, 11.3, 3.6, 3.2, 3.2, CYAN, shape=MSO_SHAPE.OVAL); grad(c2,(0x22,0xD3,0xEE),(0x6D,0x28,0xD9),120)
rect(s, -1.5, 5.3, 4.5, 4.5, VIOLET, shape=MSO_SHAPE.OVAL)
rect(s, 0.9, 1.15, 0.14, 2.4, CYAN)  # barra acento
text(s, 1.25, 1.05, 8.2, 1.2, [[R("TALLER PRÁCTICO", 18, CYAN, True, HFONT)]], space_after=0)
text(s, 1.2, 1.55, 9.5, 2.2,
     [[R("Agentes de ", 52, WHITE, True, HFONT)],
      [R("Inteligencia Artificial", 52, WHITE, True, HFONT)]],
     line_spacing=1.0, space_after=0)
text(s, 1.25, 3.7, 9.5, 0.6,
     [[R("LangChain · MCP · LangGraph · RAG · GraphRAG", 20, VIOLET_L, True, BFONT)]], space_after=0)
# chips info
def chip(x, label, val, col):
    rect(s, x, 4.7, 3.0, 1.05, INDIGO2, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    text(s, x+0.25, 4.85, 2.6, 0.4, [[R(label, 12, col, True)]], space_after=0)
    text(s, x+0.25, 5.18, 2.6, 0.5, [[R(val, 17, WHITE, True, HFONT)]], space_after=0)
chip(1.2, "DURACIÓN", "2 horas", CYAN)
chip(4.4, "ENFOQUE", "80% práctica", CORAL)
chip(7.6, "FORMATO", "Hazlo tú misma", GOLD)
footer(s, dark=True)

# ================= SLIDE 2: AGENDA (timeline) =================
s = slide(); bg(s, LIGHT)
rect(s, 0, 0, SW, 1.35, INDIGO)
rect(s, 0.6, 0.4, 0.14, 0.55, CYAN)
text(s, 0.85, 0.32, 10, 0.8, [[R("Agenda", 34, WHITE, True, HFONT)]], space_after=0)
text(s, 0.9, 0.92, 10, 0.4, [[R("Recorrido de 2 horas · demo → implementación → checkpoint", 14, RGBColor(0xB9,0xB3,0xE6))]], space_after=0)
agenda = [
    ("0:00–0:20", "LLM + Agents con LangChain", "Fundamentos y primer agente", VIOLET),
    ("0:20–0:45", "Tools con MCP", "Conectar herramientas externas", CYAN),
    ("0:45–1:10", "Migración a LangGraph", "Orquestar flujos con grafos", CORAL),
    ("1:10–1:35", "RAG", "Respuestas con documentos", GOLD),
    ("1:35–1:55", "GraphRAG", "Grafo de conocimiento", VIOLET_L),
    ("1:55–2:00", "Cierre", "Cuándo usar cada enfoque", MUTED),
]
y = 1.75
for tm, title, desc, col in agenda:
    rect(s, 0.6, y, 1.75, 0.72, col, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    text(s, 0.6, y, 1.75, 0.72, [[R(tm, 15, WHITE, True, HFONT)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, space_after=0)
    rect(s, 2.55, y, 10.15, 0.72, CARD, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    rect(s, 2.55, y, 0.12, 0.72, col)
    text(s, 2.9, y+0.07, 9.6, 0.35, [[R(title, 17, INK, True, HFONT)]], space_after=0)
    text(s, 2.9, y+0.4, 9.6, 0.3, [[R(desc, 12.5, MUTED)]], space_after=0)
    y += 0.85
footer(s); page_num(s, 2)

# ================= SLIDE 3: REGLAS =================
s = slide(); bg(s, LIGHT)
rect(s, 0, 0, 4.6, SH, INDIGO)
grad(rect(s,0,0,4.6,SH,INDIGO), (0x1A,0x16,0x4B),(0x2A,0x22,0x6E),90)
rect(s, 0.55, 1.1, 0.14, 1.9, CYAN)
text(s, 0.85, 1.0, 3.4, 2.2, [[R("Reglas del", 30, WHITE, True, HFONT)],[R("taller", 30, CYAN, True, HFONT)]], space_after=2)
text(s, 0.85, 3.2, 3.4, 2.5, [[R("Un espacio para", 14, RGBColor(0xB9,0xB3,0xE6))],[R("experimentar, equivocarse", 14, RGBColor(0xB9,0xB3,0xE6))],[R("y construir en vivo.", 14, RGBColor(0xB9,0xB3,0xE6))]], line_spacing=1.15, space_after=2)
rules = [
    ("🔁", "Ciclo de trabajo", "Demo corta → implementación guiada → checkpoint rápido", VIOLET),
    ("🙋", "Pide ayuda rápido", "Si te atoras, levanta la mano; no pierdas tiempo en silencio", CORAL),
    ("🚀", "Objetivo real", "Terminar con código corriendo en tu laptop, no solo teoría", CYAN),
]
y = 1.3
for emoji, t, d, col in rules:
    rect(s, 5.2, y, 7.5, 1.55, CARD, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    circle_icon(s, 5.55, y+0.35, 0.85, col, emoji, 26)
    text(s, 6.7, y+0.32, 5.7, 0.4, [[R(t, 19, INK, True, HFONT)]], space_after=0)
    text(s, 6.7, y+0.78, 5.7, 0.6, [[R(d, 13.5, MUTED)]], line_spacing=1.05, space_after=0)
    y += 1.8
footer(s); page_num(s, 3)

# ================= SLIDE 3b: ANTES DE EMPEZAR (SETUP) =================
s = slide(); bg(s, LIGHT)
rect(s, 0, 0, SW, 1.35, INDIGO)
rect(s, 0.6, 0.4, 0.14, 0.55, GOLD)
text(s, 0.85, 0.3, 11, 0.8, [[R("Antes de empezar", 32, WHITE, True, HFONT)]], space_after=0)
text(s, 0.9, 0.93, 11, 0.4, [[R("Deja tu entorno listo · hazlo ANTES del taller", 14, RGBColor(0xB9,0xB3,0xE6))]], space_after=0)
# comando destacado
rect(s, 0.7, 1.9, 7.4, 1.15, INDIGO2, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
rect(s, 0.7, 1.9, 0.12, 1.15, GOLD)
text(s, 1.05, 2.05, 7, 0.35, [[R("UN SOLO COMANDO EN POWERSHELL", 11.5, GOLD, True, HFONT)]], space_after=0)
text(s, 1.05, 2.4, 7, 0.55, [[R(".\\setup.ps1", 26, CYAN, True, "Consolas")]], space_after=0)
setup_steps = [
    ("1", "Clona el repositorio del taller", VIOLET),
    ("2", "Corre  .\\setup.ps1  (crea venv + instala librerías)", CORAL),
    ("3", "Pega tus credenciales de Azure en el archivo .env", CYAN),
    ("4", "Verifica con  setup_check.py  hasta ver todo ✅", VIOLET_L),
]
y = 3.35
for n, txt, col in setup_steps:
    rect(s, 0.7, y, 7.4, 0.72, CARD, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    rect(s, 0.7, y, 0.55, 0.72, col)
    text(s, 0.7, y, 0.55, 0.72, [[R(n, 18, WHITE, True, HFONT)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, space_after=0)
    text(s, 1.5, y, 6.4, 0.72, [[R(txt, 14.5, INK)]], anchor=MSO_ANCHOR.MIDDLE, space_after=0, line_spacing=1.0)
    y += 0.85
# panel lateral
rect(s, 8.5, 1.9, 4.15, 4.55, INDIGO)
grad(rect(s, 8.5, 1.9, 4.15, 4.55, INDIGO), (0x2A,0x22,0x6E),(0x1A,0x16,0x4B),115)
circle_icon(s, 10.05, 2.35, 1.05, GOLD, "⚙️", 34, INDIGO)
text(s, 8.7, 3.6, 3.75, 0.5, [[R("Requisitos", 18, CYAN, True, HFONT)]], align=PP_ALIGN.CENTER, space_after=0)
for j, req in enumerate(["Windows + PowerShell", "Python 3.10 o superior", "VS Code (Python + Jupyter)", "Credenciales de Azure OpenAI"]):
    text(s, 8.7, 4.15+j*0.5, 3.75, 0.45, [[R("• "+req, 13, RGBColor(0xC5,0xC0,0xE8))]], align=PP_ALIGN.CENTER, space_after=0)
footer(s); page_num(s)

# ---------- plantilla concepto (light) ----------
def concept_slide(n, kicker, kcolor, title, subtitle, points, emoji, big=None, biglabel=None):
    s = slide(); bg(s, LIGHT)
    # panel derecho decorativo
    rect(s, 8.9, 0, SW-8.9, SH, INDIGO)
    grad(rect(s, 8.9, 0, SW-8.9, SH, INDIGO), (0x2A,0x22,0x6E),(0x1A,0x16,0x4B),115)
    circle_icon(s, 10.4, 1.5, 1.7, kcolor, emoji, 54)
    if big:
        text(s, 8.9, 3.7, SW-8.9, 1.2, [[R(big, 60, CYAN, True, HFONT)]], align=PP_ALIGN.CENTER, space_after=0)
        text(s, 8.9, 4.9, SW-8.9, 0.8, [[R(biglabel, 14, RGBColor(0xB9,0xB3,0xE6))]], align=PP_ALIGN.CENTER, space_after=0)
    # izquierda
    rect(s, 0.7, 0.9, 0.5, 0.5, kcolor, shape=MSO_SHAPE.OVAL)
    text(s, 0.7, 0.92, 0.5, 0.5, [[R(str(n), 18, WHITE, True, HFONT)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, space_after=0)
    text(s, 1.4, 0.92, 7, 0.45, [[R(kicker, 15, kcolor, True, HFONT)]], anchor=MSO_ANCHOR.MIDDLE, space_after=0)
    text(s, 0.7, 1.55, 7.9, 1.3, [[R(title, 34, INK, True, HFONT)]], line_spacing=0.95, space_after=0)
    text(s, 0.72, 2.75, 7.9, 0.5, [[R(subtitle, 15, MUTED, italic=True)]], space_after=0)
    y = 3.55
    for pt in points:
        rect(s, 0.7, y, 7.75, 0.92, CARD, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        rect(s, 0.7, y, 0.11, 0.92, kcolor)
        text(s, 1.05, y+0.16, 7.2, 0.6, [[R(pt, 15.5, INK)]], anchor=MSO_ANCHOR.MIDDLE, space_after=0, line_spacing=1.0)
        y += 1.05
    footer(s); page_num(s)
    return s

# ---------- plantilla practica ----------
def practice_slide(num, title, time, steps, checkpoint, color):
    s = slide(); bg(s, INDIGO)
    grad(rect(s, 0, 0, SW, SH, INDIGO), (0x1A,0x16,0x4B),(0x2A,0x22,0x6E),120)
    rect(s, 9.2, -1.5, 5, 5, color, shape=MSO_SHAPE.OVAL)
    o=rect(s, 9.2, -1.5, 5, 5, color, shape=MSO_SHAPE.OVAL); grad(o,(color[0],color[1],color[2]),(0x22,0xD3,0xEE),40)
    # badge PRÁCTICA
    rect(s, 0.7, 0.75, 3.4, 0.62, color, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    text(s, 0.7, 0.75, 3.4, 0.62, [[R("PRÁCTICA "+str(num)+"  ·  MANOS A LA OBRA", 12.5, WHITE, True, HFONT)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, space_after=0)
    rect(s, 9.7, 0.8, 2.9, 0.55, None, line=color, line_w=1.5, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    text(s, 9.7, 0.8, 2.9, 0.55, [[R("⏱  "+time, 14, CYAN, True, "Segoe UI Emoji")]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, space_after=0)
    text(s, 0.7, 1.55, 11, 1.0, [[R(title, 34, WHITE, True, HFONT)]], space_after=0)
    # pasos numerados
    y = 2.95
    for i, st in enumerate(steps, 1):
        rect(s, 0.7, y, 0.62, 0.62, color, shape=MSO_SHAPE.OVAL)
        text(s, 0.7, y, 0.62, 0.62, [[R(str(i), 18, WHITE, True, HFONT)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, space_after=0)
        text(s, 1.55, y+0.02, 11, 0.6, [[R(st, 17, WHITE)]], anchor=MSO_ANCHOR.MIDDLE, space_after=0, line_spacing=1.0)
        y += 0.9
    # checkpoint
    rect(s, 0.7, 5.95, 11.95, 0.95, INDIGO2, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    rect(s, 0.7, 5.95, 0.12, 0.95, GOLD)
    circle_icon(s, 1.05, 6.2, 0.5, GOLD, "✓", 20, INDIGO)
    text(s, 1.75, 6.1, 3, 0.35, [[R("CHECKPOINT", 12, GOLD, True, HFONT)]], space_after=0)
    text(s, 1.75, 6.42, 10.6, 0.4, [[R(checkpoint, 15, WHITE, True)]], space_after=0)
    footer(s, dark=True); page_num(s, dark=True)
    return s

# mapa de numeros de pagina por titulo (para orden real)
titles_order = [
 "¿Qué es un LLM?","¿Qué es un Agent?","Práctica 1: Agente base con LangChain",
 "¿Qué es MCP?","Arquitectura MCP","Práctica 2: Agente + tools vía MCP",
 "¿Por qué migrar a LangGraph?","¿Qué es LangGraph?","Práctica 3: De Agent a Grafo",
 "¿Qué es RAG?","Componentes de RAG","Práctica 4: RAG con LangChain",
 "¿Qué es GraphRAG?","Práctica 5: Mini GraphRAG","Cuándo usar cada enfoque"]
n_map = {t: i+4 for i, t in enumerate(titles_order)}

# ================= SLIDE 4: LLM =================
concept_slide(1, "FUNDAMENTO", VIOLET, "¿Qué es un LLM?",
    "Large Language Model — Modelo de Lenguaje Grande",
    ["Modelo entrenado para comprender y generar texto",
     "Funciona mejor con contexto claro y prompts específicos",
     "Es la base para construir los agentes modernos"],
    "🧠", big="LLM", biglabel="Large Language Model")

# ================= SLIDE 5: AGENT =================
concept_slide(2, "FUNDAMENTO", CYAN, "¿Qué es un Agent?",
    "Agente — un LLM que decide y actúa",
    ["Agente = LLM + instrucciones + herramientas",
     "Flujo: observar → decidir → actuar → responder",
     "Va mucho más allá de un chatbot tradicional"],
    "🤖", big="🔁", biglabel="observar · decidir · actuar")

# ================= SLIDE 6: PRACTICA 1 =================
practice_slide(1, "Agente base con LangChain", "0:13–0:20",
    ["Crear un agente simple con un prompt template",
     "Probar el agente con 3 prompts distintos",
     "Observar cómo responde y ajustar el prompt"],
    "Todas obtienen una respuesta válida del agente", VIOLET)

# ================= SLIDE 7: MCP =================
concept_slide(3, "HERRAMIENTAS", CORAL, "¿Qué es MCP?",
    "Model Context Protocol — Protocolo de Contexto de Modelo",
    ["Protocolo para conectar modelos con herramientas externas",
     "Estandariza la integración entre herramienta y modelo",
     "Facilita la interoperabilidad y la escalabilidad"],
    "🔌", big="MCP", biglabel="Model Context Protocol")

# ================= SLIDE 8: ARQUITECTURA MCP =================
s = slide(); bg(s, LIGHT)
rect(s, 0, 0, SW, 1.35, INDIGO)
rect(s, 0.6, 0.4, 0.14, 0.55, CORAL)
text(s, 0.85, 0.3, 11, 0.8, [[R("Arquitectura MCP", 32, WHITE, True, HFONT)]], space_after=0)
text(s, 0.9, 0.93, 11, 0.4, [[R("Cómo se conectan tu agente y las herramientas", 14, RGBColor(0xB9,0xB3,0xE6))]], space_after=0)
boxes = [
    ("💻", "Cliente MCP", "Tu app o agente\nque hace las peticiones", VIOLET),
    ("🖥️", "Servidor MCP", "Expone y gestiona\nlas herramientas disponibles", CORAL),
    ("🧰", "Herramientas", "Archivos, APIs\ny servicios externos", CYAN),
]
bx = 0.85; bw = 3.7; gap = 0.55
for i,(emoji,t,d,col) in enumerate(boxes):
    x = bx + i*(bw+gap)
    rect(s, x, 2.4, bw, 3.2, CARD, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    rect(s, x, 2.4, bw, 0.18, col)
    circle_icon(s, x+bw/2-0.6, 2.95, 1.2, col, emoji, 40)
    text(s, x, 4.3, bw, 0.5, [[R(t, 21, INK, True, HFONT)]], align=PP_ALIGN.CENTER, space_after=0)
    for j,line in enumerate(d.split("\n")):
        text(s, x+0.25, 4.85+j*0.35, bw-0.5, 0.35, [[R(line, 13.5, MUTED)]], align=PP_ALIGN.CENTER, space_after=0)
    if i < 2:
        text(s, x+bw+0.02, 3.35, gap, 0.6, [[R("→", 30, VIOLET_L, True)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, space_after=0)
footer(s); page_num(s)
# ================= SLIDE 9: PRACTICA 2 =================
practice_slide(2, "Agente + tools vía MCP", "0:27–0:45",
    ["Conectar 1–2 herramientas al agente por MCP",
     "Resolver consultas que obliguen a usar las tools",
     "Revisar el rastro de llamadas a herramientas"],
    "Evidenciar que el agente usó las herramientas", CORAL)

# ================= SLIDE 10: POR QUE LANGGRAPH =================
concept_slide(4, "ORQUESTACIÓN", GOLD, "¿Por qué migrar a LangGraph?",
    "Cuando el agente necesita más control",
    ["Control explícito del flujo de ejecución",
     "Manejo de estado compartido entre pasos",
     "Ramas, reintentos y lógica condicional"],
    "🧭", big="→", biglabel="de agente a grafo")

# ================= SLIDE 11: QUE ES LANGGRAPH =================
concept_slide(5, "ORQUESTACIÓN", VIOLET_L, "¿Qué es LangGraph?",
    "Framework de grafos para orquestar agentes",
    ["Framework de grafos para orquestar agentes",
     "Nodos = pasos · Edges = transiciones",
     "Ideal para flujos complejos o de larga duración"],
    "🕸️", big="◯→◯", biglabel="nodos y transiciones")

# ================= SLIDE 12: PRACTICA 3 =================
practice_slide(3, "De Agent a Grafo", "0:54–1:10",
    ["Modelar el agente actual como un grafo",
     "Agregar un nodo de validación o clasificación",
     "Conectar los nodos con sus transiciones"],
    "Flujo de 3 o más nodos funcionando", GOLD)

# ================= SLIDE 13: RAG =================
concept_slide(6, "CONOCIMIENTO", CYAN, "¿Qué es RAG?",
    "Retrieval-Augmented Generation — Generación Aumentada por Recuperación",
    ["Recupera contexto relevante antes de responder",
     "Combina búsqueda + generación de texto",
     "Reduce alucinaciones y mejora la precisión"],
    "📚", big="RAG", biglabel="Retrieval-Augmented Generation")

# ================= SLIDE 14: COMPONENTES RAG (pipeline) =================
s = slide(); bg(s, LIGHT)
rect(s, 0, 0, SW, 1.35, INDIGO)
rect(s, 0.6, 0.4, 0.14, 0.55, CYAN)
text(s, 0.85, 0.3, 11, 0.8, [[R("Componentes de RAG", 32, WHITE, True, HFONT)]], space_after=0)
text(s, 0.9, 0.93, 11, 0.4, [[R("El pipeline que convierte documentos en respuestas", 14, RGBColor(0xB9,0xB3,0xE6))]], space_after=0)
steps = [
    ("📄","Documentos","Fuentes de\nconocimiento", VIOLET),
    ("✂️","Chunks","Fragmentos\nmanejables", CORAL),
    ("🔢","Embeddings","Vectores\nnuméricos", GOLD),
    ("🗄️","Vector store","Búsqueda por\nsimilitud", CYAN),
    ("🤖","LLM + contexto","Respuesta\nfundamentada", VIOLET_L),
]
bx=0.7; bw=2.15; gap=0.28
for i,(emoji,t,d,col) in enumerate(steps):
    x=bx+i*(bw+gap)
    rect(s, x, 2.7, bw, 2.7, CARD, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    circle_icon(s, x+bw/2-0.5, 3.05, 1.0, col, emoji, 34)
    text(s, x, 4.25, bw, 0.5, [[R(t, 15.5, INK, True, HFONT)]], align=PP_ALIGN.CENTER, space_after=0)
    for j,line in enumerate(d.split("\n")):
        text(s, x, 4.75+j*0.32, bw, 0.32, [[R(line, 11.5, MUTED)]], align=PP_ALIGN.CENTER, space_after=0)
    if i<4:
        text(s, x+bw-0.05, 3.35, gap+0.1, 0.5, [[R("→", 22, VIOLET_L, True)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, space_after=0)
footer(s); page_num(s)

# ================= SLIDE 14b: CÓMO SE MIDE LA SIMILITUD (fórmulas) =========
s = slide(); bg(s, LIGHT)
rect(s, 0, 0, SW, 1.35, INDIGO)
rect(s, 0.6, 0.4, 0.14, 0.55, CYAN)
text(s, 0.85, 0.3, 12, 0.8, [[R("¿Cómo se mide la similitud?", 32, WHITE, True, HFONT)]], space_after=0)
text(s, 0.9, 0.93, 12, 0.4, [[R("Comparar dos embeddings (vectores) para encontrar el más parecido", 14, RGBColor(0xB9,0xB3,0xE6))]], space_after=0)
PALE = RGBColor(0xEC,0xE9,0xF9)
MFONT = "Cambria Math"
formulas = [
    (CYAN,  "•", "Producto punto", "a · b = Σ aᵢ bᵢ",
     ["Multiplica y suma componente", "a componente.", "Mayor = más alineados."]),
    (VIOLET,"📏", "Distancia euclidiana", "d(a,b) = √ Σ (aᵢ − bᵢ)²",
     ["Distancia en línea recta.", "Menor = más cercanos.", "FAISS la usa por defecto (L2)."]),
    (CORAL, "∠", "Similitud coseno", "cos θ = (a · b) / (‖a‖ ‖b‖)",
     ["Mide el ángulo; ignora la", "magnitud del vector.", "1 = iguales · 0 = sin relación."]),
]
bx = 0.85; bw = 3.7; gap = 0.55
for i,(col,emoji,titulo,formula,desc) in enumerate(formulas):
    x = bx + i*(bw+gap)
    rect(s, x, 1.75, bw, 4.9, CARD, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    rect(s, x, 1.75, bw, 0.18, col)
    circle_icon(s, x+bw/2-0.55, 2.15, 1.1, col, emoji, 34)
    text(s, x, 3.4, bw, 0.5, [[R(titulo, 18, INK, True, HFONT)]], align=PP_ALIGN.CENTER, space_after=0)
    # caja de fórmula
    rect(s, x+0.3, 4.0, bw-0.6, 0.85, PALE, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    text(s, x+0.3, 4.0, bw-0.6, 0.85, [[R(formula, 16.5, INDIGO, True, MFONT)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, space_after=0)
    for j,line in enumerate(desc):
        text(s, x+0.25, 5.1+j*0.42, bw-0.5, 0.4, [[R(line, 12.5, MUTED)]], align=PP_ALIGN.CENTER, space_after=0)
text(s, 0.85, 6.8, 11.9, 0.4, [[R("a, b = embeddings · ‖a‖ = magnitud (norma) del vector · Σ = suma de todas las componentes", 11.5, MUTED, italic=True)]], align=PP_ALIGN.CENTER, space_after=0)
footer(s); page_num(s)

# ================= SLIDE 14c: BÚSQUEDA LÉXICA (TF-IDF y BM25) ==============
s = slide(); bg(s, LIGHT)
rect(s, 0, 0, SW, 1.35, INDIGO)
rect(s, 0.6, 0.4, 0.14, 0.55, GOLD)
text(s, 0.85, 0.3, 12, 0.8, [[R("Búsqueda por palabras clave: TF-IDF y BM25", 30, WHITE, True, HFONT)]], space_after=0)
text(s, 0.9, 0.93, 12, 0.4, [[R("Recuperación léxica (por términos exactos) · complementa a los embeddings", 14, RGBColor(0xB9,0xB3,0xE6))]], space_after=0)
PALE2 = RGBColor(0xEC,0xE9,0xF9)
MFONT2 = "Cambria Math"

# --- Banda 1: TF-IDF ---
rect(s, 0.7, 1.7, 11.95, 2.05, CARD, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
rect(s, 0.7, 1.7, 0.15, 2.05, GOLD)
circle_icon(s, 1.1, 2.15, 1.05, GOLD, "🔑", 32)
text(s, 2.45, 1.9, 4.0, 0.45, [[R("TF-IDF", 21, INK, True, HFONT)]], space_after=0)
text(s, 2.45, 2.36, 4.0, 0.35, [[R("Term Frequency – Inverse Document Frequency", 11.5, MUTED, italic=True)]], space_after=0)
text(s, 2.45, 2.8, 4.0, 0.8, [[R("Pesa un término por su frecuencia en el documento (tf) y su rareza en la colección (idf). Las palabras raras pesan más.", 12.5, MUTED)]], line_spacing=1.05, space_after=0)
rect(s, 6.9, 1.95, 5.5, 1.55, PALE2, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
text(s, 6.9, 1.95, 5.5, 1.55,
     [[R("tf-idf(t, d) = tf(t, d) × idf(t)", 16, INDIGO, True, MFONT2)],
      [R("idf(t) = log( N / df(t) )", 16, INDIGO, True, MFONT2)]],
     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.3, space_after=4)

# --- Banda 2: BM25 ---
rect(s, 0.7, 3.95, 11.95, 2.35, CARD, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
rect(s, 0.7, 3.95, 0.15, 2.35, VIOLET_L)
circle_icon(s, 1.1, 4.55, 1.05, VIOLET_L, "📈", 32)
text(s, 2.45, 4.15, 4.1, 0.45, [[R("BM25", 21, INK, True, HFONT)]], space_after=0)
text(s, 2.45, 4.61, 4.1, 0.35, [[R("Best Matching 25 · estándar en buscadores", 11.5, MUTED, italic=True)]], space_after=0)
text(s, 2.45, 5.05, 4.1, 1.1, [[R("Evolución de TF-IDF: satura la frecuencia (k₁) y normaliza por la longitud del documento (b). Muy usado en motores de búsqueda.", 12.5, MUTED)]], line_spacing=1.05, space_after=0)
rect(s, 6.9, 4.2, 5.5, 1.85, PALE2, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
text(s, 6.95, 4.2, 5.4, 1.85,
     [[R("score(D, Q) = Σₜ IDF(t) ·", 15, INDIGO, True, MFONT2)],
      [R("tf · (k₁ + 1)", 14.5, INDIGO, True, MFONT2)],
      [R("─────────────────────", 12, MUTED, False, MFONT2)],
      [R("tf + k₁·(1 − b + b·|D|/avgdl)", 14.5, INDIGO, True, MFONT2)]],
     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.05, space_after=2)

text(s, 0.7, 6.42, 11.95, 0.4, [[R("t = término · N = nº de documentos · df = documentos que contienen t · tf = frecuencia del término · |D| = longitud · avgdl = longitud media · k₁≈1.2–2, b≈0.75", 10.5, MUTED, italic=True)]], align=PP_ALIGN.CENTER, space_after=0)
footer(s); page_num(s)

# ================= SLIDE 15: PRACTICA 4 =================
practice_slide(4, "RAG con LangChain", "1:19–1:35",
    ["Indexar la historia de Miravalle en un vector store FAISS",
     "Preguntar al agente (p. ej. ¿quién fundó Miravalle?)",
     "Ver en el visor los fragmentos que FAISS recuperó"],
    "Respuesta con contexto citado de la historia", CYAN)

# ================= SLIDE 16: GRAPHRAG =================
concept_slide(7, "CONOCIMIENTO", VIOLET, "¿Qué es GraphRAG?",
    "RAG potenciado con un grafo de conocimiento",
    ["Combina RAG con un grafo de conocimiento",
     "Usa entidades y relaciones explícitas",
     "Útil para consultas multi-salto (multi-hop)"],
    "🕸️", big="RAG+", biglabel="grafo de conocimiento")

# ================= SLIDE 17: PRACTICA 5 =================
practice_slide(5, "Mini GraphRAG", "1:40–1:55",
    ["Extraer entidades y relaciones de Miravalle con el LLM",
     "Construir el grafo de conocimiento (se cachea en grafo.json)",
     "En el visor, cambiar el toggle a GraphRAG y comparar con RAG"],
    "Respuesta usando el grafo de conocimiento", VIOLET)

# ================= SLIDE 17b: VISOR EN VIVO (RAG vs GraphRAG) =================
s = slide(); bg(s, LIGHT)
rect(s, 0, 0, SW, 1.35, INDIGO)
rect(s, 0.6, 0.4, 0.14, 0.55, CYAN)
text(s, 0.85, 0.3, 11.5, 0.8, [[R("Visor en vivo · RAG vs GraphRAG", 30, WHITE, True, HFONT)]], space_after=0)
text(s, 0.9, 0.93, 11.5, 0.4, [[R("La misma pregunta, dos formas de recuperar · un toggle para comparar", 14, RGBColor(0xB9,0xB3,0xE6))]], space_after=0)
# toggle central
rect(s, 5.15, 1.65, 3.0, 0.62, INDIGO2, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
rect(s, 5.25, 1.72, 1.4, 0.48, CYAN, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
text(s, 5.25, 1.72, 1.4, 0.48, [[R("RAG", 13, INDIGO, True, HFONT)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, space_after=0)
text(s, 6.65, 1.72, 1.4, 0.48, [[R("GraphRAG", 12, WHITE, True, HFONT)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, space_after=0)
# dos tarjetas comparativas
cards = [
    (0.7, CYAN, "📚", "Modo RAG", "buscar_documentos",
     ["Recupera fragmentos de TEXTO", "Búsqueda semántica en FAISS", "Se resaltan en el espacio de embeddings 3D"]),
    (6.85, VIOLET, "🕸️", "Modo GraphRAG", "buscar_en_grafo",
     ["Recupera RELACIONES del grafo", "Búsqueda estructurada por entidades", "Se resaltan sobre el grafo de conocimiento"]),
]
for x, col, emoji, titulo, tool, pts in cards:
    rect(s, x, 2.5, 5.75, 4.05, CARD, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    rect(s, x, 2.5, 5.75, 0.18, col)
    circle_icon(s, x+0.35, 2.9, 0.95, col, emoji, 32)
    text(s, x+1.5, 2.95, 4.1, 0.45, [[R(titulo, 21, INK, True, HFONT)]], space_after=0)
    text(s, x+1.5, 3.42, 4.1, 0.35, [[R(tool+"( )", 13, MUTED, italic=True, font="Consolas")]], space_after=0)
    yy = 4.15
    for p in pts:
        circle_icon(s, x+0.35, yy+0.02, 0.34, col, "•", 16)
        text(s, x+0.85, yy, 4.7, 0.6, [[R(p, 13.5, INK)]], anchor=MSO_ANCHOR.MIDDLE, space_after=0, line_spacing=1.0)
        yy += 0.72
text(s, 0.7, 6.62, 11.95, 0.4, [[R("Cada modo es un agente independiente con UNA sola herramienta · no se mezclan", 13, MUTED, italic=True)]], align=PP_ALIGN.CENTER, space_after=0)
footer(s); page_num(s)

# ================= SLIDE 18: CIERRE =================
s = slide(); bg(s, INDIGO)
grad(rect(s,0,0,SW,SH,INDIGO),(0x1A,0x16,0x4B),(0x2A,0x22,0x6E),120)
rect(s, 10.5, -1.6, 4.5, 4.5, VIOLET, shape=MSO_SHAPE.OVAL)
o=rect(s,10.5,-1.6,4.5,4.5,VIOLET,shape=MSO_SHAPE.OVAL); grad(o,(0x6D,0x28,0xD9),(0x22,0xD3,0xEE),40)
rect(s, 0.7, 0.75, 0.14, 1.0, CYAN)
text(s, 1.0, 0.68, 10, 0.6, [[R("CIERRE", 15, CYAN, True, HFONT)]], space_after=0)
text(s, 0.95, 1.05, 11, 1.0, [[R("¿Cuándo usar cada enfoque?", 34, WHITE, True, HFONT)]], space_after=0)
cierre = [
    ("🧠","LLM solo","Tareas simples de texto", VIOLET_L),
    ("🔌","Agent + MCP","Cuando necesitas herramientas externas", CORAL),
    ("🧭","LangGraph","Flujos complejos con estado y ramas", GOLD),
    ("📚","RAG","Responder sobre tus documentos", CYAN),
    ("🕸️","GraphRAG","Relaciones entre entidades / multi-hop", VIOLET),
]
y=2.15
for emoji,t,d,col in cierre:
    rect(s, 0.7, y, 11.9, 0.82, INDIGO2, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    rect(s, 0.7, y, 0.12, 0.82, col)
    circle_icon(s, 1.0, y+0.16, 0.5, col, emoji, 20)
    text(s, 1.75, y, 3.2, 0.82, [[R(t, 18, WHITE, True, HFONT)]], anchor=MSO_ANCHOR.MIDDLE, space_after=0)
    text(s, 5.0, y, 7.4, 0.82, [[R(d, 15, RGBColor(0xC5,0xC0,0xE8))]], anchor=MSO_ANCHOR.MIDDLE, space_after=0)
    y+=0.95
text(s, 0.7, 6.95, 11.9, 0.4, [[R("¡Gracias! Ahora tienes 5 herramientas para construir agentes de IA 🚀", 15, CYAN, True, "Segoe UI Emoji")]], align=PP_ALIGN.CENTER, space_after=0)

import os as _os
out = _os.path.join(_os.path.dirname(_os.path.abspath(__file__)), "Taller_Agentes_IA_2h.pptx")
if _os.environ.get("DECK_OUT"):
    out = _os.environ["DECK_OUT"]
prs.save(out)
print("Guardado:", out, "· slides:", len(prs.slides._sldIdLst))
